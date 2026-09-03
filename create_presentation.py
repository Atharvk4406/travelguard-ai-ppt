from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Color Palette
NAVY = RGBColor(26, 35, 50)
IVORY = RGBColor(245, 241, 232)
BLUE = RGBColor(74, 111, 165)
TEAL = RGBColor(45, 138, 138)
AMBER = RGBColor(212, 165, 116)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(240, 238, 235)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(66)
    p.font.bold = True
    p.font.color.rgb = IVORY
    p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(36)
    p.font.color.rgb = AMBER
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, bg_color=NAVY):
    """Add a content slide with title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = TEAL
    title_shape.line.color.rgb = TEAL
    
    # Add title text
    title_frame = title_shape.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = IVORY
    p.space_before = Pt(12)
    p.space_after = Pt(12)
    title_frame.margin_left = Inches(0.5)
    
    return slide

def add_text_box(slide, left, top, width, height, text, font_size=24, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    """Add a text box to a slide"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return textbox

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # SLIDE 1: COVER
    add_title_slide(prs, "TravelGuard AI", "Predict. Detect. Recover.")
    
    # SLIDE 2: THE PROBLEM
    slide = add_content_slide(prs, "The Problem: Cascading Disruptions")
    
    # Problem description
    problem_text = "A single journey involves multiple connected legs:\n\n✈️  Flight • 🚕 Transfer • 🚆 Train • 🏨 Hotel\n\nOne delay breaks the entire itinerary."
    add_text_box(slide, 0.5, 1.2, 9, 2, problem_text, font_size=26, color=IVORY)
    
    # Example box
    add_text_box(slide, 0.8, 3.5, 4, 1.2, "Flight Arrives: 2:30 PM\nTrain Departs: 2:00 PM\nTransfer: 60 min\n\n❌ Connection Broken", 
                 font_size=18, color=AMBER, bold=True)
    
    # Current state
    add_text_box(slide, 5.5, 3.5, 3.7, 1.2, "Travelers must:\n• Check status\n• Calculate feasibility\n• Find alternatives\n• Rebuild itinerary", 
                 font_size=16, color=IVORY)
    
    # SLIDE 3: SOLUTION OVERVIEW
    slide = add_content_slide(prs, "Our Solution: Three Phases")
    
    # Phase boxes
    phases = [
        ("PREDICT", "ML models forecast\ndelay probability", 1),
        ("DETECT", "Neo4j graph identifies\nbroken connections", 4),
        ("RECOVER", "AI finds context-aware\nalternatives", 7)
    ]
    
    for i, (phase, desc, left_pos) in enumerate(phases):
        # Phase box
        phase_shape = slide.shapes.add_shape(1, Inches(left_pos), Inches(1.5), Inches(2.5), Inches(1))
        phase_shape.fill.solid()
        phase_shape.fill.fore_color.rgb = BLUE
        phase_shape.line.color.rgb = AMBER
        phase_shape.line.width = Pt(3)
        
        tf = phase_shape.text_frame
        p = tf.paragraphs[0]
        p.text = phase
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = IVORY
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(6)
        
        # Description
        add_text_box(slide, left_pos - 0.3, 2.7, 3.1, 2, desc, font_size=16, color=IVORY)
        
        # Arrow
        if i < len(phases) - 1:
            add_text_box(slide, left_pos + 2.7, 1.8, 0.6, 0.5, "→", font_size=32, color=AMBER)
    
    # SLIDE 4: TECHNICAL ARCHITECTURE
    slide = add_content_slide(prs, "Technical Architecture")
    
    # Architecture layers
    layers = [
        ("FRONTEND", "React Dashboard", 1.2),
        ("BACKEND", "Node.js + Express", 2.3),
        ("DATA LAYER", "MongoDB • Neo4j", 3.4),
        ("INTELLIGENCE", "Python + FastAPI", 4.5),
        ("AI ORCHESTRATION", "Groq LLM Agent", 5.6)
    ]
    
    for i, (layer, tech, top_pos) in enumerate(layers):
        # Layer box
        layer_shape = slide.shapes.add_shape(1, Inches(1), Inches(top_pos), Inches(8), Inches(0.7))
        layer_shape.fill.solid()
        layer_shape.fill.fore_color.rgb = BLUE if i % 2 == 0 else TEAL
        layer_shape.line.color.rgb = AMBER
        layer_shape.line.width = Pt(1)
        
        tf = layer_shape.text_frame
        p = tf.paragraphs[0]
        p.text = f"{layer}  —  {tech}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = IVORY
        p.space_before = Pt(4)
        tf.margin_left = Inches(0.3)
        
        # Arrow between layers
        if i < len(layers) - 1:
            add_text_box(slide, 4.7, top_pos + 0.75, 0.6, 0.3, "↓", font_size=24, color=AMBER)
    
    # SLIDE 5: ML DELAY PREDICTION
    slide = add_content_slide(prs, "ML Delay Prediction Engine")
    
    add_text_box(slide, 0.5, 1.2, 9, 0.6, "Historical Data → Feature Engineering → Model → Delay Probability", 
                 font_size=18, color=IVORY)
    
    # Left side - Features
    add_text_box(slide, 0.8, 2, 4, 1, "Input Features:\n• Airline/Route\n• Time of Day\n• Historical Delays\n• Weather\n• Season", 
                 font_size=16, color=IVORY)
    
    # Center - Model
    model_shape = slide.shapes.add_shape(1, Inches(3.2), Inches(3.2), Inches(3.6), Inches(1.2))
    model_shape.fill.solid()
    model_shape.fill.fore_color.rgb = TEAL
    model_shape.line.color.rgb = AMBER
    model_shape.line.width = Pt(2)
    
    tf = model_shape.text_frame
    p = tf.paragraphs[0]
    p.text = "XGBoost / Random Forest"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = IVORY
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(8)
    
    # Right side - Output
    add_text_box(slide, 5.2, 2, 4, 1.5, "Output:\n✈️ Flight AI123\nDelay Risk: HIGH\nProbability: 82%", 
                 font_size=16, color=AMBER, bold=True)
    
    add_text_box(slide, 0.8, 4.5, 8.4, 2, "Why This Matters: Early Warning\nThe traveler gets predictive alerts before actual disruption occurs, enabling proactive recovery planning.",
                 font_size=15, color=IVORY)
    
    # SLIDE 6: NEO4J TRAVEL GRAPH
    slide = add_content_slide(prs, "Neo4j Travel Graph: Dependency Modeling")
    
    add_text_box(slide, 0.5, 1.2, 9, 0.5, "Journey as a Connected Graph — Not Just Data Rows", 
                 font_size=18, bold=True, color=AMBER)
    
    # Graph visualization text
    graph_text = """Traveler → Trip → Flight AI123 → Delhi Airport → Transfer → Train 12916 → Jaipur → Hotel

Relationships modeled:
• DEPARTS_FROM | ARRIVES_AT | CONNECTS_TO | REQUIRES_TRANSFER | ALTERNATIVE_TO"""
    
    add_text_box(slide, 0.8, 2, 8.4, 2, graph_text, font_size=15, color=IVORY)
    
    add_text_box(slide, 0.8, 4.3, 8.4, 2.5, "Benefit: Cascading Disruption Detection\nWhen Flight AI123 is delayed → arrival time changes → transfer window violated → train connection becomes impossible → downstream itinerary disrupted.\n\nNeo4j naturally traverses these relationships to identify cascading failures.",
                 font_size=15, color=IVORY)
    
    # SLIDE 7: DISRUPTION DETECTION
    slide = add_content_slide(prs, "Real-Time Disruption Detection")
    
    # Real-time status
    add_text_box(slide, 0.8, 1.2, 4, 1.5, "Live Transport Status:\n\n✈️ Flight AI123\nStatus: DELAYED\nCurrent Delay: +2h 15m\nEst. Arrival: 2:30 PM", 
                 font_size=16, color=AMBER, bold=True)
    
    # Arrow
    add_text_box(slide, 4.2, 1.8, 1.6, 0.5, "Graph Analysis", font_size=14, bold=True, color=IVORY, align=PP_ALIGN.CENTER)
    
    # Connection check
    add_text_box(slide, 5.8, 1.2, 3.4, 1.5, "Connection Feasibility:\n\n🚆 Train 12916\nDeparts: 2:00 PM\nRequired Transfer: 60 min\n\n❌ NOT FEASIBLE", 
                 font_size=16, color=AMBER, bold=True)
    
    add_text_box(slide, 0.8, 3, 8.4, 3, "Disruption Engine Logic:\nThis is deterministic—not AI guesswork. The system mathematically compares:\n1. Flight expected arrival time vs. Train departure time\n2. Required transfer duration vs. available time\n3. Outputs: Connection possible or impossible\n\nOnce disruption is detected, the Recovery Engine activates.",
                 font_size=15, color=IVORY)
    
    # SLIDE 8: RECOVERY ENGINE
    slide = add_content_slide(prs, "Intelligent Recovery Engine")
    
    add_text_box(slide, 0.5, 1.2, 9, 0.5, "Multi-Criteria Ranking of Alternative Plans", 
                 font_size=18, bold=True, color=AMBER)
    
    # Recovery options
    options = [
        ("PLAN A\nAlternative Flight + Train\n₹6,400 | Low Risk\n1 Transfer | 10:45 PM", 1.5),
        ("PLAN B\nDirect Train\n₹4,800 | Medium Risk\n0 Transfers | 11:30 PM", 5),
        ("PLAN C\nOvernight + Morning Train\nHotel + 6:30 AM Train\nLowest Risk | 7:00 AM", 8.5)
    ]
    
    for i, (option, left) in enumerate(options):
        opt_shape = slide.shapes.add_shape(1, Inches(left), Inches(2.2), Inches(1.3), Inches(2.5))
        opt_shape.fill.solid()
        opt_shape.fill.fore_color.rgb = BLUE if i == 0 else LIGHT_GRAY
        opt_shape.line.color.rgb = AMBER if i == 0 else BLUE
        opt_shape.line.width = Pt(3 if i == 0 else 1)
        
        tf = opt_shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = option
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = IVORY if i == 0 else NAVY
        p.space_before = Pt(4)
        p.space_after = Pt(4)
        tf.margin_left = Inches(0.1)
    
    add_text_box(slide, 0.8, 5, 8.4, 1.8, "Ranking Criteria:\n30% Connection Reliability  •  25% Travel Time  •  20% Cost  •  15% Delay Risk  •  10% User Preference",
                 font_size=14, color=IVORY, align=PP_ALIGN.CENTER)
    
    # SLIDE 9: AGENTIC AI ORCHESTRATION
    slide = add_content_slide(prs, "Agentic AI Orchestration")
    
    add_text_box(slide, 0.5, 1.2, 9, 0.5, "AI Agent Coordinates All Systems—Not the Source of Truth", 
                 font_size=18, bold=True, color=AMBER)
    
    # Center agent
    agent_shape = slide.shapes.add_shape(1, Inches(4), Inches(2), Inches(2), Inches(1.2))
    agent_shape.fill.solid()
    agent_shape.fill.fore_color.rgb = TEAL
    agent_shape.line.color.rgb = AMBER
    agent_shape.line.width = Pt(3)
    
    tf = agent_shape.text_frame
    p = tf.paragraphs[0]
    p.text = "Recovery Agent\n(Orchestrator)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = IVORY
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(6)
    
    # Connected systems
    systems = [
        ("Transport\nStatus", 1.5, 2.4),
        ("ML Delay\nPrediction", 1.5, 3.8),
        ("Neo4j\nGraph", 7.5, 2.4),
        ("Alternative\nSearch", 7.5, 3.8)
    ]
    
    for sys, left, top in systems:
        add_text_box(slide, left, top, 1.3, 0.8, sys, font_size=12, color=IVORY, align=PP_ALIGN.CENTER)
    
    add_text_box(slide, 0.8, 5.2, 8.4, 1.8, "Agent Output: Structured Recovery Plan\nExplanation of why Plan A is recommended • Alternative options ranked • Cost/time/risk breakdown • Next steps for traveler",
                 font_size=14, color=IVORY, align=PP_ALIGN.CENTER)
    
    # SLIDE 10: DEMO SCENARIO
    slide = add_content_slide(prs, "Demo Scenario: Mumbai → Delhi → Jaipur")
    
    add_text_box(slide, 0.5, 1.1, 9, 0.4, "Real-Time Walkthrough", 
                 font_size=16, bold=True, color=AMBER)
    
    # Demo timeline
    timeline = """
T=0  Initial Plan: Flight 10:00→12:15  |  Train 14:00→18:30
      Status: 🟢 ON TRACK

T=1  ML Alert: Flight delay probability = 82%  |  Risk = HIGH
      User sees predictive warning

T=2  Real-Time Update: Flight delayed +2h 15m
      Expected arrival now 14:30

T=3  Detection: Graph analysis shows train connection broken
      ⚠️ Itinerary Disrupted

T=4  Recovery: Agent generates 3 alternatives
      Plan A (Recommended): Alternative flight + later train

T=5  User Action: Books recommended recovery plan
      New itinerary confirmed"""
    
    add_text_box(slide, 0.8, 1.7, 8.4, 4.8, timeline, font_size=13, color=IVORY)
    
    # SLIDE 11: INNOVATION & DIFFERENTIATION
    slide = add_content_slide(prs, "Innovation: What Sets Us Apart")
    
    innovations = [
        ("Delay Probability\nNot Just Alerts", "We predict risk before disruption occurs"),
        ("Dependency-Aware\nTravel Graph", "Understand how one delay cascades"),
        ("Cascading Detection\nAutomated", "Graph naturally exposes downstream failures"),
        ("Multimodal Recovery", "Flight + Train + Bus + Hotel combinations"),
        ("Context-Aware\nPlanning", "Overnight stays are smart, not failures"),
        ("Continuous Learning", "Historical outcomes improve predictions")
    ]
    
    cols = 2
    for idx, (title, desc) in enumerate(innovations):
        row = idx // cols
        col = idx % cols
        left = 0.8 + col * 4.6
        top = 1.3 + row * 1.8
        
        # Innovation box
        inn_shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(4.2), Inches(1.5))
        inn_shape.fill.solid()
        inn_shape.fill.fore_color.rgb = BLUE
        inn_shape.line.color.rgb = AMBER
        inn_shape.line.width = Pt(2)
        
        tf = inn_shape.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = AMBER
        p.space_after = Pt(4)
        tf.margin_left = Inches(0.2)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = IVORY
        p.space_before = Pt(2)
        tf.margin_left = Inches(0.2)
    
    # SLIDE 12: KEY DIFFERENTIATOR
    slide = add_content_slide(prs, "Key Differentiator")
    
    add_text_box(slide, 0.5, 1.2, 9, 1.2, "Traditional Travel Apps Follow:\nSEARCH → BOOK → TRACK", 
                 font_size=20, bold=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    
    add_text_box(slide, 0.5, 2.8, 9, 0.6, "TravelGuard Follows:", 
                 font_size=18, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
    
    add_text_box(slide, 0.5, 3.6, 9, 1.2, "PLAN → PREDICT → MONITOR → DETECT → REPLAN → RECOVER", 
                 font_size=22, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    
    add_text_box(slide, 0.8, 5.2, 8.4, 1.8, "\"TravelGuard doesn't just tell travelers that something went wrong.\nIt understands how the disruption affects their journey and helps them recover.\"",
                 font_size=16, bold=True, color=IVORY, align=PP_ALIGN.CENTER)
    
    # SLIDE 13: TECHNOLOGY STACK & DEPLOYMENT
    slide = add_content_slide(prs, "Technology Stack & Future Scope")
    
    # Tech stack
    add_text_box(slide, 0.5, 1.2, 4.5, 0.5, "Core Stack:", font_size=14, bold=True, color=AMBER)
    
    stack_text = """React
Node.js + Express
MongoDB
Neo4j
Python + FastAPI
XGBoost/Random Forest
Groq LLM
LangGraph"""
    
    add_text_box(slide, 0.8, 1.8, 4, 2.3, stack_text, font_size=13, color=IVORY)
    
    # Future scope
    add_text_box(slide, 5.2, 1.2, 4.3, 0.5, "Future Roadmap:", font_size=14, bold=True, color=AMBER)
    
    future_text = """Real airline APIs
Real railway APIs
Automated rebooking
Travel insurance integration
Weather-aware prediction
Dynamic pricing analysis
Carbon-aware recovery
Multi-city optimization"""
    
    add_text_box(slide, 5.5, 1.8, 4, 2.3, future_text, font_size=13, color=IVORY)
    
    # Deployment
    add_text_box(slide, 0.8, 4.4, 8.4, 2.6, "Deployment Strategy\nFrontend: Vercel  •  Backend: Render/Railway  •  ML Service: AWS  •  Database: MongoDB Atlas + Neo4j Aura\n\nScalability: Microservices architecture with independent scaling for ML inference, graph queries, and API servers.",
                 font_size=14, color=IVORY, align=PP_ALIGN.CENTER)
    
    # SLIDE 14: CLOSING
    slide = add_content_slide(prs, "Thank You")
    
    add_text_box(slide, 0.5, 2.2, 9, 1, "TravelGuard AI", 
                 font_size=48, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
    
    add_text_box(slide, 0.5, 3.3, 9, 0.8, "Predict. Detect. Recover.", 
                 font_size=32, color=IVORY, align=PP_ALIGN.CENTER)
    
    add_text_box(slide, 0.5, 4.4, 9, 1.5, "Questions?", 
                 font_size=28, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    
    # Save presentation
    prs.save('TravelGuard_AI_Presentation.pptx')
    print("✅ Presentation created: TravelGuard_AI_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
